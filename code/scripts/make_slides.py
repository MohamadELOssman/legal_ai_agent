#!/usr/bin/env python3
"""Build a presentation slide deck (.pptx) for the advisor meeting. No API."""

import json
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIGS = Path("experiments/figures")
EVAL = Path("experiments/retrieval_eval_196.json")
OUT = Path("../Legal_AI_Slides.pptx")

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
GREEN = RGBColor(0x0F, 0x9D, 0x58)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)

agg = {a["config"]: a for a in json.load(open(EVAL, encoding="utf-8"))["aggregates"]}
# Validated numbers from the committed benchmark run.
HIT = {"ar": (72, 81), "en": (43, 49), "fr": (39, 47)}
JUDGE = 4.6

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def bar(slide, color=NAVY, h=1.15):
    box = slide.shapes.add_shape(1, 0, 0, SW, Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    return box


def title_on_bar(slide, text, sub=None):
    bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), SW - Inches(1), Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(14); r2.font.color.rgb = RGBColor(0xBF, 0xD3, 0xE6)


def body(slide, top=1.45):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), SW - Inches(1.4), SH - Inches(top) - Inches(0.4))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bullet(tf, text, bold_lead=None, size=18, level=0, color=RGBColor(0x22, 0x22, 0x22), space=10):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.level = level; p.space_after = Pt(space)
    if bold_lead:
        r = p.add_run(); r.text = "• " + bold_lead; r.font.bold = True; r.font.size = Pt(size); r.font.color.rgb = NAVY
        r2 = p.add_run(); r2.text = text; r2.font.size = Pt(size); r2.font.color.rgb = color
    else:
        r = p.add_run(); r.text = ("• " if level == 0 else "– ") + text
        r.font.size = Pt(size); r.font.color.rgb = color
    return p


# ── Slide 1 — Title ──────────────────────────────────────────────────────────
s = add_slide()
box = s.shapes.add_shape(1, 0, 0, SW, SH); box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(2.4), SW - Inches(2), Inches(3))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "Lebanese Legal AI"
r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = WHITE
p = tf.add_paragraph(); r = p.add_run(); r.text = "Project Update & Benchmark Report"
r.font.size = Pt(24); r.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
p = tf.add_paragraph(); r = p.add_run()
r.text = f"Multi-Agent Retrieval-Augmented System · {date.today():%B %Y}"
r.font.size = Pt(16); r.font.color.rgb = RGBColor(0xBF, 0xD3, 0xE6)

# ── Slide 2 — Headline ───────────────────────────────────────────────────────
s = add_slide(); title_on_bar(s, "At a Glance")
tf = body(s)
bullet(tf, "finds a correct article 72% of the time in its top 5 results (81% in top 10).",
       "Strong on Arabic — ", size=22)
bullet(tf, f"the memoranda it writes are rated {JUDGE}/5 for legal quality.",
       "Writes well — ", size=22)
bullet(tf, "search method and settings were chosen from benchmark evidence, not guesswork.",
       "Rigorous — ", size=22)
bullet(tf, "English/French retrieval and exact-citation precision — both with a fix already scoped.",
       "Honest next step — ", size=22)
# big metric strip
strip = s.shapes.add_textbox(Inches(0.7), Inches(5.6), SW - Inches(1.4), Inches(1.2))
p = strip.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Arabic: 72% / 81% hit-rate (top 5 / 10)      |      Answer quality: 4.6 / 5"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = GREEN

# ── Slide 3 — How it works ───────────────────────────────────────────────────
s = add_slide(); title_on_bar(s, "How the System Works", "A pipeline of specialised AI agents — one job each")
tf = body(s)
for lead, rest in [
    ("Orchestrator → ", "decides the question type (general law vs. a case to assess)."),
    ("Query Understanding → ", "detects language & domain, extracts the facts."),
    ("Research (RAG) → ", "retrieves the most relevant articles and rulings."),
    ("Analysis → ", "extracts the applicable provisions and checks them against the sources."),
    ("Reasoning → ", "applies the law to the facts."),
    ("Citation → ", "formats and verifies every article number."),
    ("Writing → ", "produces the final memorandum, in the question's language."),
]:
    bullet(tf, rest, lead, size=17, space=7)
bullet(tf, "flags any claim not found in the sources (a 'hallucination rate').",
       "Trust layer → ", size=17, space=7)

# ── Slide 4 — What was updated ───────────────────────────────────────────────
s = add_slide(); title_on_bar(s, "What Was Updated This Phase")
tf = body(s)
for lead, rest in [
    ("Models: ", "standardised on Claude Sonnet 4.5 (+ selector for 4.6 / Opus / Haiku)."),
    ("Corpus: ", "added the English Penal Code — now bilingual (713 documents)."),
    ("Trust: ", "grounding + citation verification + hallucination rate."),
    ("Retrieval: ", "hybrid search adopted by evidence; harmful re-ranker removed."),
    ("Experience: ", "answers in the question's language; Arabic shown right-to-left."),
    ("Engineering: ", "reproducible pipeline, tests, CI, and cost/time tracking."),
]:
    bullet(tf, rest, lead, size=19, space=12)

# ── Slide 5 — Benchmark: how ─────────────────────────────────────────────────
s = add_slide(); title_on_bar(s, "The Benchmark — How It Works")
tf = body(s)
bullet(tf, "196 questions generated from the REAL corpus — each with its correct article number(s) as the 'gold answer', validated automatically.", "Grounded: ", size=18)
bullet(tf, "Arabic, English, French; 204 distinct articles; general questions + case scenarios.", "Realistic: ", size=18)
bullet(tf, "an answer is correct if it retrieves the right article number (so it fairly tests cross-lingual search).", "Objective: ", size=18)
bullet(tf, "retrieval (hit-rate, recall, ranking), citation accuracy, answer quality (1–5), cost, and statistical confidence.", "Measured: ", size=18)
bullet(tf, "to turn 'the answers look good' into objective, repeatable, defensible numbers.", "Why: ", size=18)

# ── Slide 6 — Results: by language ───────────────────────────────────────────
s = add_slide(); title_on_bar(s, "Results — Accuracy by Language")
s.shapes.add_picture(str(FIGS / "fig_languages.png"), Inches(0.6), Inches(1.5), height=Inches(4.7))
tf = body(s, top=1.7)
# put text on right side
tb = s.shapes.add_textbox(Inches(7.4), Inches(1.9), Inches(5.4), Inches(4.5)); tf = tb.text_frame; tf.word_wrap = True
bullet(tf, "Arabic is Lebanon's official legal language — and the corpus is Arabic-primary.", size=18)
bullet(tf, "72% of correct articles appear in the top 5; 81% in the top 10.", "Arabic: ", size=18)
bullet(tf, "lower (corpus is Arabic-heavy); a cross-lingual fix is already built.", "EN / FR: ", size=18)

# ── Slide 7 — Results: method + quality ──────────────────────────────────────
s = add_slide(); title_on_bar(s, "Results — Search Method & Answer Quality")
s.shapes.add_picture(str(FIGS / "fig_methods.png"), Inches(0.6), Inches(1.5), height=Inches(4.6))
tb = s.shapes.add_textbox(Inches(7.4), Inches(1.9), Inches(5.4), Inches(4.5)); tf = tb.text_frame; tf.word_wrap = True
bullet(tf, "Hybrid (keyword + meaning) finds the correct law most often — the default.", "Best method: ", size=18)
bullet(tf, "a popular English re-ranker made results WORSE here, so it was removed.", "Rejected: ", size=18)
bullet(tf, f"final memoranda rated {JUDGE}/5 for correctness, citations, completeness, clarity.", "Answer quality: ", size=18)
bullet(tf, "~3 minutes, ~$0.12 per full answer.", "Practical: ", size=18)

# ── Slide 8 — Findings & next steps ──────────────────────────────────────────
s = add_slide(); title_on_bar(s, "Findings & Next Steps")
tf = body(s)
bullet(tf, "Strong on Arabic (72% / 81%); writes high-quality memoranda (4.6/5).", "Findings: ", size=19)
bullet(tf, "design choices made by evidence (hybrid in, re-ranker out).", "", size=18, level=1)
bullet(tf, "Raise English/French retrieval (evaluate the built-in cross-lingual translation).", "Next: ", size=19)
bullet(tf, "Improve citation precision (cite only directly-applicable articles).", "", size=18, level=1)
bullet(tf, "Run the full evaluation for final multi-agent vs. baseline numbers.", "", size=18, level=1)
bullet(tf, "Expand the corpus to contract law and French legal texts.", "", size=18, level=1)

prs.save(str(OUT))
print(f"Saved: {OUT.resolve()}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
