#!/usr/bin/env python3
"""Detailed academic presentation (.pptx) for the Lebanese Legal AI thesis. No API."""

from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIGS = Path("experiments/figures")
OUT = Path("../Legal_AI_Academic_Presentation.pptx")

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x0F, 0x9D, 0x58)
AMBER = RGBColor(0xB4, 0x7A, 0x00)
GREY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7)
CODEBG = RGBColor(0x0F, 0x17, 0x2A)
CODEFG = RGBColor(0xD7, 0xE3, 0xF4)
BORDER = RGBColor(0xCB, 0xD5, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUB = RGBColor(0xBF, 0xD3, 0xE6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.color.rgb = line if line else color
    if not line:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h); tb.text_frame.word_wrap = True
    return tb.text_frame


def run(p, text, size, color=GREY, bold=False, italic=False, mono=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic
    if mono:
        r.font.name = "Consolas"
    return r


def title_bar(s, title, sub=None):
    rect(s, 0, 0, SW, Inches(1.05), NAVY)
    tf = textbox(s, Inches(0.5), Inches(0.08), SW - Inches(1), Inches(0.9))
    run(tf.paragraphs[0], title, 27, WHITE, bold=True)
    if sub:
        run(tf.add_paragraph(), sub, 13, SUB)


def body(s, top=1.3):
    return textbox(s, Inches(0.7), Inches(top), SW - Inches(1.4), SH - Inches(top) - Inches(0.35))


def bullet(tf, text, lead=None, size=17, level=0, space=8, color=GREY, first=False):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.level = level; p.space_after = Pt(space)
    prefix = "• " if level == 0 else "– "
    if lead:
        run(p, prefix + lead, size, NAVY, bold=True)
        run(p, text, size, color)
    else:
        run(p, prefix + text, size, color)
    return p


def promptbox(s, x, y, w, h, title, lines, size=11):
    box = rect(s, x, y, w, h, CODEBG)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.12); tf.margin_right = Inches(0.15)
    if title:
        run(tf.paragraphs[0], title, size + 1, RGBColor(0x7A, 0xD1, 0xA0), bold=True, mono=True)
        start = tf.add_paragraph()
    else:
        start = tf.paragraphs[0]
    for i, ln in enumerate(lines):
        p = start if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run(p, ln, size, CODEFG, mono=True)
    return box


def section(num, title, subtitle=None):
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, Inches(0.9), Inches(2.9), Inches(0.16), Inches(1.5), GREEN)
    tf = textbox(s, Inches(1.3), Inches(2.75), SW - Inches(2), Inches(2.0))
    run(tf.paragraphs[0], f"Part {num}", 18, SUB, bold=True)
    run(tf.add_paragraph(), title, 34, WHITE, bold=True)
    if subtitle:
        run(tf.add_paragraph(), subtitle, 15, SUB, italic=True)
    return s


# ── 1. Title ─────────────────────────────────────────────────────────────────
s = slide(); rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.5), SW, Inches(0.08), GREEN)
tf = textbox(s, Inches(1), Inches(2.1), SW - Inches(2), Inches(2.7))
run(tf.paragraphs[0], "A Multi-Agent, Retrieval-Augmented AI System", 34, WHITE, bold=True)
run(tf.add_paragraph(), "for Lebanese Legal Research", 34, WHITE, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(16)
run(p, "Trilingual (Arabic · English · French) · Grounded · Rigorously Benchmarked", 18, SUB)
run(tf.add_paragraph(), f"Thesis Presentation · {date.today():%B %Y}", 14, SUB, italic=True)

# ── 2. Motivation ────────────────────────────────────────────────────────────
s = slide(); title_bar(s, "1. Motivation & Problem", "Why an AI system for Lebanese law")
tf = body(s)
bullet(tf, "official language is Arabic, but practice mixes French and English; sources are scattered and hard to search.",
       "Lebanese law is trilingual — ", first=True)
bullet(tf, "ordinary citizens cannot easily understand their rights or obligations.",
       "Access to justice — ")
bullet(tf, "lawyers spend hours locating the correct articles and relevant precedents for a case.",
       "Professional burden — ")
bullet(tf, "generic chatbots invent article numbers and cannot be trusted for legal citation (hallucination).",
       "Reliability gap — ")
bullet(tf, "answer legal questions by RETRIEVING the correct law, then writing a grounded, cited answer whose accuracy can be measured.",
       "Goal — ")

# ── 3. Objectives ────────────────────────────────────────────────────────────
s = slide(); title_bar(s, "2. Research Objectives")
tf = body(s)
for lead, rest in [
    ("Multi-agent design: ", "can specialised agents (retrieval, analysis, reasoning, writing) outperform a single LLM?"),
    ("Trilingual retrieval: ", "how well can the system find the right law across Arabic / English / French?"),
    ("Trustworthiness: ", "can every legal claim and citation be grounded and verified against the corpus?"),
    ("Adaptivity: ", "can the output adapt to the user — citizen, lawyer, or judge?"),
    ("Rigorous evaluation: ", "build a benchmark that proves accuracy objectively, against a source of truth."),
]:
    bullet(tf, rest, lead, size=18, space=13, first=lead.startswith("Multi"))

# ── 4. System overview ───────────────────────────────────────────────────────
s = slide(); title_bar(s, "3. System Overview")
tf = body(s)
bullet(tf, "A pipeline of specialised AI agents turns a legal question into a grounded, cited answer.",
       "What it does: ", size=18, first=True)
bullet(tf, "Retrieval-Augmented Generation (RAG) over the Lebanese Penal Code + court rulings.",
       "Approach: ", size=18)
bullet(tf, "Arabic, English, and French; the answer is written in the question's language.",
       "Languages: ", size=18)
bullet(tf, "every article the answer relies on is checked against the corpus; a hallucination rate is reported.",
       "Trust: ", size=18)
bullet(tf, "a Streamlit web app + a rigorous, reproducible benchmark to measure quality.",
       "Delivered as: ", size=18)

# ── 5. Architecture diagram ──────────────────────────────────────────────────
s = slide(); title_bar(s, "4. Architecture — The Agent Pipeline")
agents = [("0", "Orchestrator"), ("1", "Query\nUnderstanding"), ("2", "Research\n(RAG)"),
          ("3", "Analysis"), ("4", "Reasoning"), ("5", "Citation"), ("6", "Writing")]
n = len(agents); gap = Inches(0.12)
bw = (SW - Inches(1.0) - gap * (n - 1)) / n
x = Inches(0.5); y = Inches(2.35); bh = Inches(1.45)
for i, (num, name) in enumerate(agents):
    col = GREEN if i == 0 else (AMBER if i == 6 else BLUE)
    box = rect(s, x, y, bw, bh, col)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, num, 20, WHITE, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    run(p2, name, 11.5, WHITE, bold=True)
    x = x + bw + gap
cap = textbox(s, Inches(0.7), Inches(4.15), SW - Inches(1.4), Inches(2.9))
bullet(cap, "classifies the question type AND the user type; routes the whole pipeline.", "Orchestrator → ", size=14.5, first=True, space=5)
bullet(cap, "detects language / domain / facts as validated structured data.", "Query Understanding → ", size=14.5, space=5)
bullet(cap, "retrieves the most relevant articles and rulings (hybrid search).", "Research → ", size=14.5, space=5)
bullet(cap, "extracts applicable provisions, grounded against the retrieved text.", "Analysis → ", size=14.5, space=5)
bullet(cap, "applies the law to the facts (chain-of-thought).", "Reasoning → ", size=14.5, space=5)
bullet(cap, "formats + verifies citations, then writes the final answer in the user's shape & language.", "Citation · Writing → ", size=14.5, space=5)

# ── 6. System prompts / agent instruction ────────────────────────────────────
s = slide(); title_bar(s, "5. How the Agents Are Instructed (System Prompts)",
                       "Each agent has a role-specific system prompt that fixes its behaviour")
tf = textbox(s, Inches(0.6), Inches(1.25), Inches(6.4), Inches(5.8))
bullet(tf, "each agent gets a dedicated system prompt defining its single job and its rules.", "Role-specific: ", size=15, first=True, space=8)
bullet(tf, "entry agents return VALIDATED structured objects (tool-use), not free text — no fragile parsing.", "Structured: ", size=15, space=8)
bullet(tf, "every generative prompt forbids inventing law: 'cite ONLY articles supplied to you'.", "Anti-hallucination: ", size=15, space=8)
bullet(tf, "the writer is told to answer entirely in the query's language, with role-specific sections.", "Language-locked: ", size=15, space=8)
bullet(tf, "prompts are stored in files and versioned, so they can be tuned and compared.", "Externalised: ", size=15, space=8)
promptbox(s, Inches(7.2), Inches(1.4), Inches(5.5), Inches(2.35), "Orchestrator (system prompt)", [
    "Read the user input, classify it, return a", "routing decision. You never answer directly.",
    "Classify TWO things:",
    "  A) query_type: general | case_analysis",
    "  B) user_type : citizen | lawyer | judge",
])
promptbox(s, Inches(7.2), Inches(3.95), Inches(5.5), Inches(2.05), "Writing / Analysis (rules)", [
    "Only cite articles supplied to you;",
    "never invent legal references.",
    "Write the ENTIRE memorandum in the",
    "query's language ONLY.",
])

# ── 7. Structured outputs & grounding ────────────────────────────────────────
s = slide(); title_bar(s, "6. Structured Outputs & Grounding")
tf = textbox(s, Inches(0.6), Inches(1.25), Inches(6.6), Inches(5.8))
bullet(tf, "Orchestrator & Query-Understanding emit typed objects via tool-use (Pydantic schema) — the model MUST return valid fields.", "Schema-validated: ", size=15.5, first=True, space=9)
bullet(tf, "the Analysis agent tags each extracted provision as grounded / ungrounded by matching its article number to the retrieved text.", "Grounding check: ", size=15.5, space=9)
bullet(tf, "the Citation agent verifies each article against a master corpus index and drops loosely-related neighbours.", "Citation filter: ", size=15.5, space=9)
bullet(tf, "the Writer receives a CLOSED set — it may cite only those verified numbers.", "Closed set: ", size=15.5, space=9)
promptbox(s, Inches(7.4), Inches(1.5), Inches(5.3), Inches(2.0), "Routing schema (tool-use)", [
    "class RoutingDecision:",
    "  query_type: str   # general | case",
    "  user_type : str   # citizen|lawyer|judge",
    "  detected_language, legal_domain,",
    "  extracted_facts[], pipeline_config",
])
promptbox(s, Inches(7.4), Inches(3.7), Inches(5.3), Inches(1.7), "Closed-set constraint (writer)", [
    "CITATION CONSTRAINT: You may cite ONLY",
    "these article numbers: 547, 549.",
    "Do NOT mention any other article.",
])

# ── 8. Adaptive output (user types) ──────────────────────────────────────────
s = slide(); title_bar(s, "7. Adaptive Output — Who Is Asking?",
                       "The Orchestrator detects the user type and shapes the answer accordingly")
tf = body(s, top=1.5)
bullet(tf, "→ a plain, jargon-free answer to their question.", "👤 Citizen (general question) ", size=18, space=15, first=True)
bullet(tf, "→ a structured advisory memorandum (defence-oriented) for a client's case.", "⚖️ Lawyer ", size=18, space=15)
bullet(tf, "→ a formal judicial DECISION: facts → applicable law → reasoning → verdict.", "👨‍⚖️ Judge (gives facts, wants ruling) ", size=18, space=15)
bullet(tf, "The user can select the role, or the Orchestrator infers it from the phrasing (e.g. 'my client' → lawyer).", "", size=15, color=GREY)

# ── 9. Corpus ────────────────────────────────────────────────────────────────
s = slide(); title_bar(s, "8. Corpus & Data Foundation")
tf = body(s)
bullet(tf, "Lebanese Penal Code — 417 Arabic + 242 English articles + 54 Court of Cassation rulings = 713 indexed documents.",
       "Bilingual corpus: ", size=17, first=True)
bullet(tf, "a reproducible pipeline reads the sources, embeds them, and records exactly what was indexed.",
       "Ingestion: ", size=17)
bullet(tf, "new sources (contract law, French texts) are auto-indexed by the same pipeline.",
       "Extensible: ", size=17)
bullet(tf, "every article keeps a validated number — later used as the 'gold answer' in the benchmark.",
       "Provenance: ", size=17)

# ── 10. Retrieval ────────────────────────────────────────────────────────────
s = slide(); title_bar(s, "9. Retrieval (RAG)")
tf = body(s)
bullet(tf, "Hybrid search = keyword (BM25) + meaning (dense embeddings) — chosen by benchmark evidence.",
       "Method: ", size=17, first=True)
bullet(tf, "multilingual sentence-transformer (mpnet), running locally and free.",
       "Embeddings: ", size=17)
bullet(tf, "articles and rulings are searched in separate pools; the article number makes matches language-agnostic.",
       "Cross-lingual: ", size=17)
bullet(tf, "a popular English re-ranker was TESTED and REJECTED — it hurt this Arabic/legal corpus.",
       "Evidence over intuition: ", size=17)

# ── 11. Trust ────────────────────────────────────────────────────────────────
s = slide(); title_bar(s, "10. Trust & Anti-Hallucination")
tf = body(s)
bullet(tf, "each provision is checked against the retrieved text; ungrounded ones are flagged.",
       "Grounding: ", size=17, first=True)
bullet(tf, "every cited article number is verified against the corpus index (no invented citations).",
       "Citation verification: ", size=17)
bullet(tf, "the memo may cite ONLY verified articles it was given (precision-first).",
       "Closed citation set: ", size=17)
bullet(tf, "each answer reports a hallucination rate and a grounding rate.",
       "Trust metric: ", size=17)

# ── 12. Web app overview ─────────────────────────────────────────────────────
s = slide(); title_bar(s, "11. The Web Application (UI)", "Built with Streamlit — three working areas")
tf = body(s, top=1.5)
bullet(tf, "ask a question, pick your role, run all 7 agents; see the memo, sources, trust indicators, tokens/cost.",
       "🔗 End-to-End Pipeline: ", size=17, space=12, first=True)
bullet(tf, "run and inspect each agent in isolation (inputs, outputs, retrieved chunks) for debugging.",
       "🔬 Individual Agents: ", size=17, space=12)
bullet(tf, "generate test questions, enter reference answers, and score the system against a source of truth.",
       "📊 Benchmarking: ", size=17, space=12)
bullet(tf, "step-by-step transparency — every agent's output, timing, and cost are shown live.",
       "Design principle: ", size=16, space=8, color=GREY)

# ── 13. UI details ───────────────────────────────────────────────────────────
s = slide(); title_bar(s, "11. UI — Pipeline View in Detail")
tf = body(s)
for lead, rest in [
    ("Role selector: ", "Citizen / Lawyer / Judge / Auto-detect — changes the output shape."),
    ("Per-agent panels: ", "each of the 7 steps shows a ✅, its time, and its output (expandable)."),
    ("Grounding & citations: ", "provisions marked grounded/ungrounded; citations shown ✓ verified / ⚠ unverified."),
    ("Memorandum: ", "rendered as a clean document — right-to-left for Arabic, in the query's language."),
    ("Cost & trust summary: ", "tokens, USD cost, latency, grounding rate, and hallucination rate per run."),
    ("Download: ", "full results (JSON) and the memorandum are downloadable."),
]:
    bullet(tf, rest, lead, size=16.5, space=9, first=lead.startswith("Role"))

# ══ PART II — BENCHMARKING ════════════════════════════════════════════════════
section("II", "Benchmarking", "The core scientific contribution")

# ── 15. Benchmark: why ───────────────────────────────────────────────────────
s = slide(); title_bar(s, "12. Benchmark — Why It Matters")
tf = body(s)
bullet(tf, "a legal AI is only credible if its accuracy can be MEASURED, not just claimed.", "Motivation: ", size=17.5, first=True, space=11)
bullet(tf, "prove whether the multi-agent design beats a single-LLM baseline.", "Compare: ", size=17.5, space=11)
bullet(tf, "quantify how often the system finds the right law and cites the right articles.", "Quantify: ", size=17.5, space=11)
bullet(tf, "turn subjective impressions into objective, repeatable, defensible numbers.", "Rigour: ", size=17.5, space=11)
bullet(tf, "test every design decision (search method, embedding, re-ranker) with data.", "Guide design: ", size=17.5, space=11)

# ── 16. Question generation ──────────────────────────────────────────────────
s = slide(); title_bar(s, "13. Generating Benchmark Questions",
                       "A dedicated generation agent builds grounded questions from the real corpus")
tf = textbox(s, Inches(0.6), Inches(1.25), Inches(6.7), Inches(5.9))
bullet(tf, "sample real Penal-Code articles (and rulings) from the corpus.", "Step 1 — ", size=15, first=True, space=7)
bullet(tf, "the LLM writes realistic questions answerable from THOSE articles, in a target language.", "Step 2 — ", size=15, space=7)
bullet(tf, "it returns each question WITH the article numbers it is based on (the gold answer).", "Step 3 — ", size=15, space=7)
bullet(tf, "gold numbers are validated against the corpus index; invalid ones are dropped.", "Step 4 — ", size=15, space=7)
bullet(tf, "batched (~10 questions / call) for efficiency; deduped; mixes general + case questions.", "Step 5 — ", size=15, space=7)
bullet(tf, "questions are grounded in real law, so each has a known, checkable correct answer.", "Result: ", size=15, space=7, color=NAVY)
promptbox(s, Inches(7.5), Inches(1.5), Inches(5.25), Inches(3.4), "Generation prompt (excerpt)", [
    "From these Lebanese Penal Code articles:",
    "  - Article 549: ...",
    "  - Article 547: ...",
    "",
    "Generate N diverse questions in {language}.",
    "Set gold_articles to those article numbers",
    "(choose only from: 547, 548, 549).",
    "Mix general questions and case scenarios.",
])

# ── 17. What is measured (metrics table) ─────────────────────────────────────
s = slide(); title_bar(s, "14. Benchmark — What Is Measured")
tbl_rows = [
    ("Retrieval", "Precision@k, Recall@k, MRR, nDCG, hit-rate", "Does it find the correct article?"),
    ("Citations", "Precision, Recall, F1 vs. gold articles", "Does the answer cite the right articles?"),
    ("Answer quality", "LLM-as-judge 1–5 (correctness, citations, completeness, clarity)", "Is the final answer sound & well written?"),
    ("Reference-based", "Judge compares answer vs. a human source-of-truth answer", "Does it match the expert answer?"),
    ("Efficiency", "Latency, tokens, cost per query", "Is it practical to run?"),
    ("Statistics", "Mean ± 95% CI, paired significance tests", "Are differences real, not chance?"),
]
tb = s.shapes.add_table(len(tbl_rows) + 1, 3, Inches(0.6), Inches(1.35),
                        SW - Inches(1.2), Inches(5.4)).table
tb.columns[0].width = Inches(2.4); tb.columns[1].width = Inches(6.2); tb.columns[2].width = Inches(3.5)
for j, h in enumerate(["Layer", "Metrics", "Question it answers"]):
    c = tb.cell(0, j); c.text = h
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    rr = c.text_frame.paragraphs[0].runs[0]; rr.font.color.rgb = WHITE; rr.font.bold = True; rr.font.size = Pt(14)
for i, (a, b, cc) in enumerate(tbl_rows, 1):
    for j, val in enumerate((a, b, cc)):
        c = tb.cell(i, j); c.text = val
        r0 = c.text_frame.paragraphs[0].runs[0]; r0.font.size = Pt(12.5); r0.font.color.rgb = GREY
        if j == 0:
            r0.font.bold = True; r0.font.color.rgb = NAVY

# ── 18. The evaluator agent (judge) ──────────────────────────────────────────
s = slide(); title_bar(s, "15. The Evaluator Agent (LLM-as-Judge)",
                       "A separate LLM scores each answer — an automated expert reviewer")
tf = textbox(s, Inches(0.6), Inches(1.25), Inches(6.7), Inches(5.9))
bullet(tf, "a second, independent LLM (Claude), separate from the answering agents.", "What it is: ", size=15, first=True, space=8)
bullet(tf, "temperature 0 for consistent, deterministic scoring.", "Deterministic: ", size=15, space=8)
bullet(tf, "scores 4 dimensions 1–5: legal correctness, citation quality, completeness, clarity.", "Rubric: ", size=15, space=8)
bullet(tf, "returns a strict JSON verdict + a one-line explanation; the average is the score.", "Output: ", size=15, space=8)
bullet(tf, "objective citation-F1 (vs gold articles) is ALSO computed automatically to cross-check the judge.", "Cross-checked: ", size=15, space=8)
promptbox(s, Inches(7.5), Inches(1.5), Inches(5.25), Inches(4.6), "Judge prompt (reference-free)", [
    "You are a Lebanese legal evaluation expert.",
    "User Query: \"...\"",
    "Legal Memorandum: \"...\"",
    "",
    "Score each 1-5:",
    "  legal_correctness, citation_quality,",
    "  completeness, clarity",
    "Return ONLY JSON:",
    "  {legal_correctness:N, ... ,",
    "   explanation:\"one sentence\"}",
])

# ── 19. Source of truth (reference answer) ───────────────────────────────────
s = slide(); title_bar(s, "16. The Source-of-Truth (Reference) Answer",
                       "Why the human ground-truth answer is essential")
tf = textbox(s, Inches(0.6), Inches(1.25), Inches(6.7), Inches(5.9))
bullet(tf, "without a reference, the judge scores on its own opinion — which can be biased or wrong on Lebanese law.", "The problem: ", size=15, first=True, space=9)
bullet(tf, "the user enters the correct expert answer per question — the SOURCE OF TRUTH.", "The fix: ", size=15, space=9)
bullet(tf, "the judge then compares the AI answer AGAINST this reference (agreement on law, citations, completeness).", "Comparison: ", size=15, space=9)
bullet(tf, "grounds evaluation in human expertise — objective and reproducible, not the judge's guess.", "Why it matters: ", size=15, space=9)
bullet(tf, "it is REQUIRED for a judged comparison run — the run is blocked until every question has one.", "Mandatory: ", size=15, space=9)
promptbox(s, Inches(7.5), Inches(1.5), Inches(5.25), Inches(4.3), "Judge prompt (reference-based)", [
    "Compare the AI answer against the",
    "REFERENCE (ground-truth) answer.",
    "",
    "REFERENCE answer: \"...\"",
    "AI answer: \"...\"",
    "",
    "Score 1-5, judged AGAINST the reference:",
    "  correctness, citations, completeness,",
    "  clarity.",
])

# ── 20. Benchmark workflow in the app ────────────────────────────────────────
s = slide(); title_bar(s, "17. Benchmark — The Workflow in the App")
tf = body(s)
bullet(tf, "set the number of questions; the generator produces grounded questions, shown 10 per page.", "1 · Generate: ", size=16.5, first=True, space=9)
bullet(tf, "review them, then type the source-of-truth answer in the editable table (required).", "2 · Reference: ", size=16.5, space=9)
bullet(tf, "choose systems — Multi-Agent vs Single-Agent vs No-RAG.", "3 · Select: ", size=16.5, space=9)
bullet(tf, "each system answers; the judge scores every answer against the reference.", "4 · Run: ", size=16.5, space=9)
bullet(tf, "per-system scores, per-dimension table, charts, citation metrics, and downloadable JSON.", "5 · Results: ", size=16.5, space=9)

# ── 21. Results by language ──────────────────────────────────────────────────
s = slide(); title_bar(s, "18. Results — Accuracy by Language")
if (FIGS / "fig_languages.png").exists():
    s.shapes.add_picture(str(FIGS / "fig_languages.png"), Inches(0.6), Inches(1.4), height=Inches(4.7))
tf = textbox(s, Inches(7.5), Inches(1.8), Inches(5.3), Inches(4.7))
bullet(tf, "72% of correct articles in the top 5, 81% in the top 10.", "Arabic (primary legal language): ", size=17, first=True, space=13)
bullet(tf, "lower — the corpus is Arabic-primary; closing this is the main next step.", "English / French: ", size=17, space=13)
bullet(tf, "the article-number metric fairly credits cross-lingual matches.", "Note: ", size=15, color=GREY)

# ── 22. Results method + quality ─────────────────────────────────────────────
s = slide(); title_bar(s, "19. Results — Search Method & Answer Quality")
if (FIGS / "fig_methods.png").exists():
    s.shapes.add_picture(str(FIGS / "fig_methods.png"), Inches(0.6), Inches(1.4), height=Inches(4.6))
tf = textbox(s, Inches(7.5), Inches(1.8), Inches(5.3), Inches(4.7))
bullet(tf, "hybrid (keyword + meaning) finds the correct law most often — chosen default.", "Best method: ", size=17, first=True, space=13)
bullet(tf, "a popular English re-ranker made results worse and was removed.", "Rejected by data: ", size=17, space=13)
bullet(tf, "final memoranda rated ~4.6/5 for legal quality by the judge.", "Answer quality: ", size=17, space=13)
bullet(tf, "~2–3 minutes, ~$0.12 per full answer.", "Practical: ", size=17, space=13)

# ── 23. Findings ─────────────────────────────────────────────────────────────
s = slide(); title_bar(s, "20. Key Findings")
tf = body(s)
bullet(tf, "specialised agents produce grounded, well-structured answers (~4.6/5).", "Multi-agent works: ", size=17, first=True)
bullet(tf, "hybrid > semantic; the re-ranker hurts; the embedding is validated — all decided by the benchmark.", "Evidence-based engineering: ", size=17)
bullet(tf, "answer quality is bounded by retrieval recall — if the article isn't found, it can't be cited.", "Bottleneck identified: ", size=17)
bullet(tf, "tightening citations raised citation-F1 ~3× (0.04 → 0.13) on the same questions.", "Precision fix, measured: ", size=17)

# ── 24. Engineering ──────────────────────────────────────────────────────────
s = slide(); title_bar(s, "21. Engineering & Reproducibility")
tf = body(s)
for lead, rest in [
    ("Standardised model: ", "Claude Sonnet 4.5 (selectable: 4.6 / Opus / Haiku)."),
    ("Headless pipeline: ", "the full system runs without the UI for batch evaluation."),
    ("Tests + CI: ", "unit tests and GitHub Actions run automatically on every change."),
    ("Telemetry: ", "tokens, cost, and latency tracked per agent."),
    ("Reproducible: ", "one command rebuilds the index; the benchmark is regenerable."),
    ("Version-controlled: ", "all work committed and pushed to GitHub."),
]:
    bullet(tf, rest, lead, size=16.5, space=9, first=lead.startswith("Stand"))

# ── 25. Limitations ──────────────────────────────────────────────────────────
s = slide(); title_bar(s, "22. Limitations & Future Work")
tf = body(s)
bullet(tf, "English/French retrieval lags Arabic (~40% vs 72%).", "Limitation: ", size=17, first=True)
bullet(tf, "corpus is criminal (Penal) law only; single-article gold is strict for multi-article questions.", "Scope: ", size=17)
bullet(tf, "evaluate the built-in cross-lingual translation; test stronger multilingual embeddings.", "Next — retrieval: ", size=17)
bullet(tf, "add contract law + French texts (pipeline supports it); collect expert reference answers at scale.", "Next — corpus & rigor: ", size=17)
bullet(tf, "structured reasoning, a verifier agent, and a full statistical comparison run.", "Next — agents & eval: ", size=17)

# ── 26. Conclusion ───────────────────────────────────────────────────────────
s = slide(); rect(s, 0, 0, SW, SH, NAVY); rect(s, 0, Inches(1.15), SW, Inches(0.06), GREEN)
tf = textbox(s, Inches(0.8), Inches(0.35), SW - Inches(1.6), Inches(0.8))
run(tf.paragraphs[0], "23. Conclusion & Contributions", 28, WHITE, bold=True)
tf = textbox(s, Inches(0.9), Inches(1.55), SW - Inches(1.8), Inches(5.5))
for lead, rest in [
    ("A working trilingual legal AI ", "for Lebanese law: grounded, cited, adaptive to the user."),
    ("A rigorous benchmark ", "of grounded questions with validated gold answers and multi-metric scoring."),
    ("Reference-based evaluation ", "— answers judged against a human source of truth, not the judge's opinion."),
    ("Evidence-based design ", "— every retrieval choice proven (or rejected) by measurement."),
    ("Trust by construction ", "— grounding, citation verification, and a reported hallucination rate."),
    ("Reproducible & engineered ", "— tests, CI, telemetry, and a usable web application."),
]:
    p = tf.add_paragraph(); p.space_after = Pt(12)
    run(p, "✓ " + lead, 17, GREEN, bold=True); run(p, rest, 17, WHITE)

prs.save(str(OUT))
print(f"Saved: {OUT.resolve()}  ({len(prs.slides._sldIdLst)} slides)")
